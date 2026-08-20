# -*- coding: utf-8 -*-
"""上传文件格式转换（MarkItDown 优先，独立解析器回退）。

统一把可文本化的文件转换为 **Markdown (.md)** 文本。与历史版本不同，
本模块不再直接落盘——转换在 **host 侧**（第三方库）完成，返回 markdown
文本字符串，由上层（routers/uploads.py）经 ``workspace.get_backend()``
写入沙箱；该设计使上传逻辑沙箱感知（双 PVC / 共享 PVC 下 session 隔离）。

转换策略（双通道）：
- 文档/HTML 类（PDF / Word(docx) / PPT(pptx) / Excel(xlsx,xlsm) / HTML）**优先**
  由 ``markitdown`` 转换（输出质量更高，表格/标题识别好）；MarkItDown 失败
  （格式不支持、转换异常、输出为空）时**回退**到该格式的独立解析器：
  pdfplumber（pdf）/ python-docx（docx）/ python-pptx（pptx）/
  openpyxl（xlsx,xlsm）/ html2text（html）。
- 老格式 .doc/.ppt/.xls 双通道（MarkItDown 与独立解析器）均不支持，已从
  支持列表移除——上传后作为"仅原始文件"处理，不尝试转换。
- 文本/代码类（txt/md/csv/json/xml/log/各类源码）为复制语义，直接按
  UTF-8 解码返回（不走转换器，避免源码类扩展名不支持导致失败）。
- 图片、压缩包、二进制等 -> 由调用方按需拒绝。

依赖（pyproject.toml service extra）：``markitdown[all]`` + 上述
独立解析器库（回退通道需要）。
"""
from __future__ import annotations

import io

from .manager import UploadError


# 文本/代码类：复制语义（非转换），UTF-8 直读
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl", ".xml",
    ".log", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".sh", ".bash",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".go", ".c", ".cpp", ".h",
    ".hpp", ".cs", ".rb", ".php", ".rs", ".kt", ".swift", ".sql", ".r", ".scala",
    ".pl", ".lua", ".vim", ".dockerfile", ".gitignore", ".env",
}

# 文档/HTML 类：值为入库的 format 标签（MarkItDown 优先 + 独立解析器回退）
# 仅保留双通道（MarkItDown 与独立解析器）都支持的格式；老格式 .doc/.ppt/.xls
# 双通道均不支持，已移除——上传后作为"仅原始文件"处理，不再尝试转换。
_DOC_FORMAT = {
    ".pdf": "pdf",
    ".docx": "word",
    ".pptx": "ppt",
    ".xlsx": "excel",
    ".xlsm": "excel",
    ".html": "html",
    ".htm": "html",
}


class UnsupportedFileType(UploadError):
    """文件类型不在支持范围内。"""


def is_supported_format(filename: str) -> bool:
    """判断文件名是否可转换为 .md。"""
    ext = f".{_split_ext(filename)}"  # 补点后与下方带点扩展名集合比较
    return ext in _TEXT_EXTS or ext in _DOC_FORMAT


def _split_ext(filename: str) -> str:
    name = (filename or "").strip().replace("\\", "/").rsplit("/", 1)[-1]
    return name.lower().rsplit(".", 1)[-1] if "." in name else ""


def convert_file_bytes(
    filename: str,
    data: bytes,
    content_type: str | None = None,
) -> tuple[str, str]:
    """把上传文件的字节内容转换为 markdown 文本（host 侧执行）。

    Args:
        filename: 客户端原始文件名（含扩展名）。
        data: 已读取的文件字节。
        content_type: 可选 MIME 类型（当前仅作候选判定，未强制）。

    Returns:
        ``(format, markdown_text)``；无法转换时抛
        `UnsupportedFileType` 或具体转换错误（调用方捕获后不阻断上传）。
    """
    ext = f".{_split_ext(filename)}"
    if ext in _TEXT_EXTS:
        # 文本/代码类：复制语义，UTF-8 直读（非转换）
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("utf-8", errors="ignore")
        return "text", text

    if ext in _DOC_FORMAT:
        fmt = _DOC_FORMAT[ext]
        # 优先 MarkItDown；失败（不支持/异常/空输出）回退该格式独立解析器
        try:
            return fmt, _convert_markitdown(data, ext)
        except UnsupportedFileType:
            fallback = _FALLBACK_CONVERTERS.get(ext)
            if fallback is None:
                raise
            return fmt, fallback(data)

    raise UnsupportedFileType(f"unsupported file type: {ext}")


def _convert_markitdown(src: bytes, ext: str) -> str:
    """用 MarkItDown 把文档字节转换为 markdown 文本。

    按扩展名交给 MarkItDown 内置转换器（pdf/docx/pptx/xlsx/html）；
    任何失败（未安装、格式不支持、转换异常、输出为空）都抛
    `UnsupportedFileType`，由调用方回退到独立解析器。
    """
    try:
        from markitdown import MarkItDown
    except ImportError as e:
        raise UnsupportedFileType("markitdown not installed") from e
    try:
        result = MarkItDown().convert_stream(
            io.BytesIO(src),
            file_extension=ext,
        )
    except Exception as e:  # noqa: BLE001
        raise UnsupportedFileType(f"markitdown failed for {ext}: {e}") from e
    text = (result.text_content or "").strip()
    if not text:
        raise UnsupportedFileType(f"markitdown produced empty output for {ext}")
    return text


# ---------------------------------------------------------------------------
# 独立解析器回退通道（MarkItDown 失败时按扩展名选用）
# 统一签名：(src: bytes) -> md_text: str；库缺失时抛 UnsupportedFileType
# ---------------------------------------------------------------------------
_FALLBACK_CONVERTERS: dict[str, object] = {}


def _convert_pdf(src: bytes) -> str:
    try:
        import pdfplumber  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("pdfplumber not installed") from e
    parts: list[str] = []
    with pdfplumber.open(io.BytesIO(src)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            parts.append(f"## 第 {i} 页\n\n{text}")
    return "\n\n".join(parts)


def _convert_docx(src: bytes) -> str:
    try:
        from docx import Document  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("python-docx not installed") from e
    doc = Document(io.BytesIO(src))
    lines: list[str] = []
    for para in doc.paragraphs:
        style = (para.style.name or "") if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading"):
            level = "".join(filter(str.isdigit, style)) or "1"
            lines.append(f"{'#' * min(int(level), 6)} {text}")
        else:
            lines.append(text)
    return "\n\n".join(lines)


def _convert_pptx(src: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("python-pptx not installed") from e
    prs = Presentation(io.BytesIO(src))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        lines.append(t)
    return "\n\n".join(lines)


def _convert_excel(src: bytes) -> str:
    try:
        import openpyxl  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("openpyxl not installed") from e
    wb = openpyxl.load_workbook(io.BytesIO(src), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"## Sheet: {ws.title}")
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        header = [str(c) if c is not None else "" for c in rows[0]]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    return "\n".join(lines)


def _convert_html(src: bytes) -> str:
    try:
        import html2text  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("html2text not installed") from e
    h = html2text.HTML2Text()
    h.body_width = 0  # 不自动换行
    raw = src.decode("utf-8", errors="ignore")
    return h.handle(raw)


_FALLBACK_CONVERTERS.update({
    ".pdf": _convert_pdf,
    ".docx": _convert_docx,
    ".pptx": _convert_pptx,
    ".xlsx": _convert_excel,
    ".xlsm": _convert_excel,
    ".html": _convert_html,
    ".htm": _convert_html,
})
